"""
PowerPoint Generator
"""

import os
from pptx import Presentation


def create_presentation(title, content, output_folder="output"):
    """
    Create a PowerPoint presentation with one title slide.
    """

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = content

    file_path = os.path.join(output_folder, "presentation.pptx")

    prs.save(file_path)

    return file_path